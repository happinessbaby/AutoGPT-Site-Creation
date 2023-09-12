# import requests
import os
import fitz
from pathlib import Path
import pypandoc
import uuid
import markdown
import csv
import bs4
import urllib.request
from urllib.request import Request, urlopen
import uuid
from pptx import Presentation
from langchain.document_loaders import UnstructuredURLLoader
from typing import Any, List, Union, Dict



def convert_to_txt(file, output_path):
    file_ext = os.path.splitext(file)[1]
    if (file_ext)=='.txt':
        os.rename(file, output_path)
    if (file_ext=='.pdf'): 
        convert_pdf_to_txt(file, output_path)
    elif (file_ext=='.odt' or file_ext=='.docx'):
        convert_doc_to_txt(file, output_path)
    elif (file_ext==".log"):
        convert_log_to_txt(file, output_path)
    elif (file_ext==".pptx"):
        convert_pptx_to_txt(file, output_path)
    elif (file_ext==".ipynb"):
        convert_ipynb_to_txt(file, output_path)

def convert_ipynb_to_txt(file, output_path):
    os.rename(file, output_path)

        

def convert_log_to_txt(file, output_path):
    with open(file, "r") as f:
        content = f.read()
        print(content)
        with open(output_path, "w") as f:
            f.write(content)
            f.close()

def convert_pptx_to_txt(pptx_file, output_path):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text+=shape.text+'\n'
    with open(output_path, 'w') as f:
        f.write(text)
        f.close()


def convert_pdf_to_txt(pdf_file, output_path):
    pdf = fitz.open(pdf_file)
    text = ""
    for page in pdf:
        text+=page.get_text() + '\n'
    with open(output_path, 'w') as f:
        f.write(text)
        f.close()


def convert_doc_to_txt(doc_file, output_path):
    pypandoc.convert_file(doc_file, 'plain', outputfile=output_path)

def read_txt(file):
    try:
        with open(file, 'r', errors='ignore') as f:
            text = f.read()
            return text
    except Exception as e:
        print(e)
    
def markdown_table_to_dict(markdown_table):
    # Convert Markdown to HTML
    html = markdown.markdown(markdown_table)
    print(html)

    # Parse HTML table using csv module
    rows = csv.reader(html.split('\n'), delimiter='|')

    # Extract header and data rows
    header = next(rows)
    data = [row for row in rows if len(row) > 1]

    # Convert rows to dictionary
    result = []
    for row in data:
        item = {}
        for i, value in enumerate(row):
            key = header[i].strip()
            item[key] = value.strip()
        result.append(item)

    return result

class AppURLopener(urllib.request.FancyURLopener):
    version = "Mozilla/5.0"

def retrieve_web_content(link, save_path="./web_data/test.txt"):

    req = Request(
        url=link, 
        headers={'User-Agent': 'Mozilla/5.0'}
    )
    try: 
        webpage=str(urllib.request.urlopen(link).read())
    except Exception: 
        # webpage = urlopen(req).read()
        opener = AppURLopener()
        webpage = opener.open(link)
    soup = bs4.BeautifulSoup(webpage, features="lxml")

    content = soup.get_text()

    if content:
        with open(save_path, 'w') as file:
            file.write(content)
            file.close()
            print('Content retrieved and written to file.')
            return True
    else:
        print('Failed to retrieve content from the URL.')
        return False
    

def write_to_docx(doc: Any, text: str, type: str, res_path: str):
    if type=="cover_letter":
        p = doc.add_paragraph()
        p.add_run(text)
        doc.save(res_path)


        




if __name__=="__main__":
    retrieve_web_content(
        "https://learning.shine.com/talenteconomy/resume-tips/resume-objectives/",
        save_path = f"./web_data/{str(uuid.uuid4())}.txt")





    

    




